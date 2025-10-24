from flask import Flask, render_template, request, jsonify, send_file
import os
import json
from werkzeug.utils import secure_filename
import requests
from pathlib import Path
import PyPDF2
import docx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
import chromadb
from chromadb.utils import embedding_functions
import uuid
import re

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['OUTPUT_FOLDER'] = 'outputs'
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024

os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)

chroma_client = chromadb.Client()
collection = None

OLLAMA_URL = "http://localhost:11434/api/generate"
MODEL_NAME = "llama3.2:3b"

sessions = {}

# Professional color schemes
COLOR_SCHEMES = {
    'corporate_blue': {
        'primary': RGBColor(0, 51, 102),      # Dark blue
        'secondary': RGBColor(0, 102, 204),   # Blue
        'accent': RGBColor(255, 153, 0),      # Orange
        'text': RGBColor(51, 51, 51),         # Dark gray
        'light': RGBColor(240, 248, 255)      # Light blue
    },
    'modern_green': {
        'primary': RGBColor(34, 139, 34),     # Forest green
        'secondary': RGBColor(50, 205, 50),   # Lime green
        'accent': RGBColor(255, 215, 0),      # Gold
        'text': RGBColor(33, 33, 33),
        'light': RGBColor(240, 255, 240)
    },
    'elegant_purple': {
        'primary': RGBColor(75, 0, 130),      # Indigo
        'secondary': RGBColor(147, 112, 219), # Medium purple
        'accent': RGBColor(255, 192, 203),    # Pink
        'text': RGBColor(40, 40, 40),
        'light': RGBColor(248, 248, 255)
    }
}

def extract_text_from_pdf(file_path):
    text = ""
    with open(file_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
    return text

def extract_text_from_docx(file_path):
    doc = docx.Document(file_path)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def extract_text_from_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

def chunk_text(text, chunk_size=500, overlap=50):
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size - overlap):
        chunk = ' '.join(words[i:i + chunk_size])
        chunks.append(chunk)
    return chunks

def query_ollama(prompt, context=""):
    full_prompt = f"{context}\n\n{prompt}" if context else prompt
    
    payload = {
        "model": MODEL_NAME,
        "prompt": full_prompt,
        "stream": False,
        "temperature": 0.7
    }
    
    try:
        response = requests.post(OLLAMA_URL, json=payload)
        response.raise_for_status()
        return response.json()['response']
    except Exception as e:
        return f"Error querying Ollama: {str(e)}"

def apply_text_formatting(text_frame, text, font_size=18, bold=False, color=None):
    """Apply sophisticated text formatting"""
    text_frame.text = text
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.name = 'Calibri'
        paragraph.font.bold = bold
        if color:
            paragraph.font.color.rgb = color

def add_styled_title_slide(prs, title, subtitle, color_scheme):
    """Create a sophisticated title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background shape
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = color_scheme['light']
    background.line.fill.background()
    
    # Add accent bar
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width, Inches(1.5)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = color_scheme['primary']
    accent_bar.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        prs.slide_width - Inches(1), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = 'Calibri'
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Add subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2),
            prs.slide_width - Inches(1), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.name = 'Calibri'
        subtitle_frame.paragraphs[0].font.color.rgb = color_scheme['text']

def add_content_slide(prs, title, content, color_scheme, layout_type='bullet'):
    """Create sophisticated content slides with various layouts"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width, Inches(0.8)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = color_scheme['primary']
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        prs.slide_width - Inches(1), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Content based on layout type
    if layout_type == 'bullet':
        add_bullet_content(slide, content, color_scheme)
    elif layout_type == 'two_column':
        add_two_column_content(slide, content, color_scheme)
    elif layout_type == 'numbered':
        add_numbered_content(slide, content, color_scheme)

def add_bullet_content(slide, points, color_scheme):
    """Add bullet point content"""
    content_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.2),
        Inches(8.6), Inches(5.5)
    )
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.level = 0
        p.font.size = Pt(20)
        p.font.name = 'Calibri'
        p.font.color.rgb = color_scheme['text']
        p.space_after = Pt(12)
        
        # Add bullet
        p.bullet = True

def add_two_column_content(slide, content, color_scheme):
    """Add two-column layout content"""
    mid_point = len(content) // 2
    left_content = content[:mid_point]
    right_content = content[mid_point:]
    
    # Left column
    left_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.2),
        Inches(4), Inches(5.5)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, point in enumerate(left_content):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = color_scheme['text']
        p.bullet = True
    
    # Right column
    right_box = slide.shapes.add_textbox(
        Inches(5.2), Inches(1.2),
        Inches(4), Inches(5.5)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for i, point in enumerate(right_content):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = color_scheme['text']
        p.bullet = True

def add_numbered_content(slide, points, color_scheme):
    """Add numbered list content"""
    content_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.2),
        Inches(8.6), Inches(5.5)
    )
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = f"{i+1}. {point}"
        p.font.size = Pt(20)
        p.font.name = 'Calibri'
        p.font.color.rgb = color_scheme['text']
        p.space_after = Pt(12)

def add_chart_slide(prs, title, chart_data, color_scheme):
    """Add a slide with a chart"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background and header (same as content slide)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width, Inches(0.8)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = color_scheme['primary']
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        prs.slide_width - Inches(1), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Add chart
    chart_type_map = {
        'bar': XL_CHART_TYPE.BAR_CLUSTERED,
        'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'line': XL_CHART_TYPE.LINE,
        'pie': XL_CHART_TYPE.PIE
    }
    
    chart_type = chart_type_map.get(chart_data.get('type', 'column'), XL_CHART_TYPE.COLUMN_CLUSTERED)
    
    # Prepare chart data
    chart_data_obj = CategoryChartData()
    chart_data_obj.categories = chart_data.get('categories', ['Category 1', 'Category 2', 'Category 3'])
    
    for series in chart_data.get('series', [{'name': 'Series 1', 'values': [10, 20, 30]}]):
        chart_data_obj.add_series(series['name'], series['values'])
    
    # Add chart to slide
    x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(5)
    chart = slide.shapes.add_chart(
        chart_type, x, y, cx, cy, chart_data_obj
    ).chart
    
    # Style the chart
    chart.has_legend = True
    chart.legend.position = 2  # Right
    chart.legend.font.size = Pt(12)

def add_table_slide(prs, title, table_data, color_scheme):
    """Add a slide with a table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background and header
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width, Inches(0.8)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = color_scheme['primary']
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        prs.slide_width - Inches(1), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Add table
    rows = len(table_data.get('rows', []))
    cols = len(table_data.get('headers', []))
    
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(0.5) * (rows + 1)
    
    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table
    
    # Set headers
    for i, header in enumerate(table_data.get('headers', [])):
        cell = table.rows[0].cells[i]
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = color_scheme['primary']
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
    
    # Set data
    for row_idx, row_data in enumerate(table_data.get('rows', [])):
        for col_idx, cell_value in enumerate(row_data):
            cell = table.rows[row_idx + 1].cells[col_idx]
            cell.text = str(cell_value)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            
            # Alternate row colors
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = color_scheme['light']

def create_presentation(presentation_data):
    """Create sophisticated PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Select color scheme
    color_scheme_name = presentation_data.get('color_scheme', 'corporate_blue')
    color_scheme = COLOR_SCHEMES.get(color_scheme_name, COLOR_SCHEMES['corporate_blue'])
    
    # Title slide
    add_styled_title_slide(
        prs,
        presentation_data.get('title', 'Presentation'),
        presentation_data.get('subtitle', ''),
        color_scheme
    )
    
    # Content slides
    for slide_data in presentation_data.get('slides', []):
        slide_type = slide_data.get('type', 'bullet')
        
        if slide_type == 'chart':
            add_chart_slide(
                prs,
                slide_data.get('title', ''),
                slide_data.get('chart_data', {}),
                color_scheme
            )
        elif slide_type == 'table':
            add_table_slide(
                prs,
                slide_data.get('title', ''),
                slide_data.get('table_data', {}),
                color_scheme
            )
        else:
            layout_type = slide_data.get('layout', 'bullet')
            add_content_slide(
                prs,
                slide_data.get('title', ''),
                slide_data.get('points', []),
                color_scheme,
                layout_type
            )
    
    return prs

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_documents():
    global collection
    
    if 'files' not in request.files:
        return jsonify({'error': 'No files provided'}), 400
    
    files = request.files.getlist('files')
    
    if not files:
        return jsonify({'error': 'No files selected'}), 400
    
    try:
        chroma_client.delete_collection(name="documents")
    except:
        pass
    
    collection = chroma_client.create_collection(
        name="documents",
        metadata={"hnsw:space": "cosine"}
    )
    
    all_chunks = []
    chunk_ids = []
    
    for file in files:
        if file.filename == '':
            continue
        
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        
        if filename.endswith('.pdf'):
            text = extract_text_from_pdf(file_path)
        elif filename.endswith('.docx'):
            text = extract_text_from_docx(file_path)
        elif filename.endswith('.txt'):
            text = extract_text_from_txt(file_path)
        else:
            continue
        
        chunks = chunk_text(text)
        all_chunks.extend(chunks)
        chunk_ids.extend([f"{filename}_{i}" for i in range(len(chunks))])
    
    if all_chunks:
        collection.add(
            documents=all_chunks,
            ids=chunk_ids
        )
    
    return jsonify({
        'success': True,
        'message': f'Processed {len(files)} files with {len(all_chunks)} chunks'
    })

@app.route('/generate_presentation', methods=['POST'])
def generate_presentation():
    data = request.json
    user_request = data.get('request', '')
    session_id = data.get('session_id', str(uuid.uuid4()))
    
    if not user_request:
        return jsonify({'error': 'No request provided'}), 400
    
    if session_id not in sessions:
        sessions[session_id] = {
            'history': [],
            'iterations': 0
        }
    
    context = ""
    if collection:
        results = collection.query(
            query_texts=[user_request],
            n_results=5
        )
        if results['documents']:
            context = "\n\n".join(results['documents'][0])
    
    previous_context = ""
    if sessions[session_id]['history']:
        previous_context = "\n\nPrevious presentation attempts:\n"
        for i, hist in enumerate(sessions[session_id]['history'], 1):
            previous_context += f"\nAttempt {i}:\n{json.dumps(hist['structure'], indent=2)}\n"
            previous_context += f"User feedback: {hist['feedback']}\n"
    
    prompt = f"""Based on the following context and user request, create a SOPHISTICATED PowerPoint presentation structure.

Context from documents:
{context}

{previous_context}

User request: {user_request}

Create a JSON structure with the following features:
- Use various slide types: "bullet", "two_column", "numbered", "chart", "table"
- Include professional layouts
- For chart slides, provide chart data with categories and series
- For table slides, provide headers and rows
- Choose a color scheme: "corporate_blue", "modern_green", or "elegant_purple"

JSON format:
{{
    "title": "Main presentation title",
    "subtitle": "Subtitle",
    "color_scheme": "corporate_blue",
    "slides": [
        {{
            "type": "bullet",
            "layout": "bullet",
            "title": "Slide title",
            "points": ["Point 1", "Point 2", "Point 3"]
        }},
        {{
            "type": "chart",
            "title": "Chart title",
            "chart_data": {{
                "type": "column",
                "categories": ["Q1", "Q2", "Q3"],
                "series": [
                    {{"name": "Sales", "values": [100, 150, 200]}}
                ]
            }}
        }},
        {{
            "type": "table",
            "title": "Table title",
            "table_data": {{
                "headers": ["Column 1", "Column 2"],
                "rows": [["Data 1", "Data 2"], ["Data 3", "Data 4"]]
            }}
        }}
    ]
}}

Provide ONLY valid JSON, no additional text."""
    
    response = query_ollama(prompt, context)
    
    try:
        start_idx = response.find('{')
        end_idx = response.rfind('}') + 1
        json_str = response[start_idx:end_idx]
        presentation_structure = json.loads(json_str)
    except:
        presentation_structure = {
            "title": "Generated Presentation",
            "subtitle": "Based on your request",
            "color_scheme": "corporate_blue",
            "slides": [
                {
                    "type": "bullet",
                    "layout": "bullet",
                    "title": "Overview",
                    "points": ["Point 1", "Point 2", "Point 3"]
                }
            ]
        }
    
    sessions[session_id]['iterations'] += 1
    
    return jsonify({
        'session_id': session_id,
        'structure': presentation_structure,
        'iteration': sessions[session_id]['iterations']
    })

@app.route('/confirm_presentation', methods=['POST'])
def confirm_presentation():
    data = request.json
    confirmed = data.get('confirmed', False)
    session_id = data.get('session_id')
    structure = data.get('structure')
    feedback = data.get('feedback', '')
    
    if not session_id or session_id not in sessions:
        return jsonify({'error': 'Invalid session'}), 400
    
    if confirmed:
        try:
            prs = create_presentation(structure)
            output_filename = f"presentation_{session_id}.pptx"
            output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
            prs.save(output_path)
            
            return jsonify({
                'success': True,
                'message': 'Presentation generated successfully',
                'download_url': f'/download/{output_filename}'
            })
        except Exception as e:
            return jsonify({'error': f'Error generating presentation: {str(e)}'}), 500
    else:
        sessions[session_id]['history'].append({
            'structure': structure,
            'feedback': feedback
        })
        
        return jsonify({
            'success': True,
            'message': 'Feedback recorded. Please submit a new generation request.'
        })

@app.route('/download/<filename>')
def download_file(filename):
    file_path = os.path.join(app.config['OUTPUT_FOLDER'], filename)
    if os.path.exists(file_path):
        return send_file(file_path, as_attachment=True)
    return jsonify({'error': 'File not found'}), 404

if __name__ == '__main__':
    app.run(debug=True, port=5000)
